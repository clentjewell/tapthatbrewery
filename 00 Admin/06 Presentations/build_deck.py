"""JP_TapThat_WhereYouAre_v01 — adapted from the jp-brand-presentation reference generator."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

JEWELL_BLACK = RGBColor(0x11,0x11,0x11); CREAM = RGBColor(0xFA,0xF8,0xF4)
SIGNAL_BLUE  = RGBColor(0x2D,0x5B,0xFF); CALM_GREY = RGBColor(0xED,0xEB,0xEA)
GREY_TEXT    = RGBColor(0x66,0x66,0x66)
POPPINS="Poppins"; SLIDE_W=Inches(13.333); SLIDE_H=Inches(7.5); MARGIN=Inches(1.0)

def bg(slide,c):
    f=slide.background.fill; f.solid(); f.fore_color.rgb=c

def tb(slide,left,top,width,height,text,size_pt=14,bold=False,color=JEWELL_BLACK,
       uppercase=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,line_spacing=1.5,tracking=None):
    box=slide.shapes.add_textbox(left,top,width,height); tf=box.text_frame
    tf.word_wrap=True; tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    tf.vertical_anchor=anchor
    for i,line in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.line_spacing=line_spacing
        r=p.add_run(); r.text=line.upper() if uppercase else line
        f=r.font; f.name=POPPINS; f.size=Pt(size_pt); f.bold=bold; f.color.rgb=color
        if tracking is not None: r._r.get_or_add_rPr().set("spc",str(tracking))
    return box

def hairline(slide,left,top,width,color=CALM_GREY,weight_pt=0.75):
    ln=slide.shapes.add_connector(1,left,top,left+width,top)
    ln.line.color.rgb=color; ln.line.width=Pt(weight_pt); return ln

def s_title(prs,bar,title,subtitle=None,date_line=None):
    s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s,CREAM)
    tb(s,MARGIN,Inches(0.4),Inches(11.3),Inches(0.3),bar,9,uppercase=True,tracking=200)
    tb(s,MARGIN,SLIDE_H-MARGIN-Inches(2.6),Inches(11.3),Inches(2.0),title,60,bold=True,
       line_spacing=1.05,anchor=MSO_ANCHOR.BOTTOM)
    if subtitle: tb(s,MARGIN,SLIDE_H-Inches(1.35),Inches(11.3),Inches(0.6),subtitle,20)
    if date_line: tb(s,Inches(9.0),SLIDE_H-Inches(0.6),Inches(3.3),Inches(0.3),date_line,11,
                     color=GREY_TEXT,align=PP_ALIGN.RIGHT)
    return s

def s_content(prs,eyebrow,headline,body=None,head_size=38,body_top=4.5,body_w=10.0):
    s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s,CREAM)
    tb(s,MARGIN,MARGIN,Inches(11.3),Inches(0.3),eyebrow,9,uppercase=True,tracking=200)
    tb(s,MARGIN,Inches(1.9),Inches(11.3),Inches(2.3),headline,head_size,line_spacing=1.12)
    if body: tb(s,MARGIN,Inches(body_top),Inches(body_w),Inches(2.4),body,14,line_spacing=1.55)
    return s

def s_divider(prs,num,title,strapline=None):
    s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s,JEWELL_BLACK)
    tb(s,MARGIN,MARGIN,Inches(11.3),Inches(0.3),num,12,color=CREAM,uppercase=True,tracking=200)
    tb(s,MARGIN,Inches(2.3),Inches(11.3),Inches(3.0),title,96,color=CREAM,
       line_spacing=1.0,anchor=MSO_ANCHOR.MIDDLE)
    if strapline: tb(s,MARGIN,Inches(5.9),Inches(11.3),Inches(0.6),strapline,16,color=CREAM)
    return s

def s_closer(prs,action,owner):
    s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s,CREAM)
    tb(s,MARGIN,Inches(2.0),Inches(11.3),Inches(2.0),"Next.",96,bold=True,line_spacing=1.0)
    tb(s,MARGIN,Inches(4.5),Inches(11.3),Inches(0.5),action,14)
    tb(s,MARGIN,Inches(5.0),Inches(11.3),Inches(0.5),owner,14)
    return s

def stat_row(slide,items,top=4.4):
    """Three or four figures with a label. Type and space, no chart furniture."""
    n=len(items); w=11.3/n
    for i,(num,lbl) in enumerate(items):
        L=MARGIN+Inches(w*i)
        hairline(slide,L,Inches(top),Inches(w-0.4))
        tb(slide,L,Inches(top+0.18),Inches(w-0.4),Inches(0.7),num,32,bold=True,line_spacing=1.0)
        tb(slide,L,Inches(top+0.95),Inches(w-0.5),Inches(0.9),lbl,11,color=GREY_TEXT,line_spacing=1.35)

def rule_table(slide,rows,top=2.6,col=4.4,size=13,step=0.78):
    """Rows separated by hairlines. Replaces bullets, which the deck standard bans."""
    y=top; h=step-0.06
    for k,v in rows:
        hairline(slide,MARGIN,Inches(y),Inches(11.3))
        tb(slide,MARGIN,Inches(y+0.14),Inches(col-0.3),Inches(h),k,size,bold=True,line_spacing=1.3)
        tb(slide,MARGIN+Inches(col),Inches(y+0.14),Inches(11.3-col),Inches(h),v,size,
           color=GREY_TEXT,line_spacing=1.38)
        y+=step
    hairline(slide,MARGIN,Inches(y),Inches(11.3))
    assert y<=7.05, f"rule_table overruns the slide: last rule at {y:.2f}in"

prs=Presentation(); prs.slide_width=SLIDE_W; prs.slide_height=SLIDE_H
BAR="JP · TAP THAT BREWERY · WHERE YOU ARE"

# 01
s_title(prs,BAR,"Where you are, and the\nthree things to do about it.",
        "Prepared for Justin Mistry and Harry. The detail sits behind every slide.",
        "28 August 2026")

# 02
s_divider(prs,"Part 01","Where you are","We have looked at this properly. This is the read.")

# 03
s=s_content(prs,"The read",
  "A refill business is being run with a\nhospitality venue's attention. And it is\nlosing money.",
  None,head_size=36)
stat_row(s,[("206","Systems in your database"),
            ("96 to 113","Active customers. Your own documents define active three ways"),
            ("56%","Bought their system somewhere else and refill with you anyway"),
            ("~9x","The climb to 1,000 from where you actually are")])

# 04
s_content(prs,"The problem underneath",
  "The taproom and the refill model\nwork against each other.",
  "The refill business is built on people not coming in. So every campaign that fills the "
  "taproom fills the half that loses money. We had the taproom down as a funnel. It is at least "
  "as much a cost, and nobody has priced what carrying 27 taps actually costs you.")

# 05
s=s_content(prs,"Before anything else","What we got wrong.",None)
rule_table(s,[("We never said it","Nothing in seventy nine documents said the business is losing money."),
              ("The taproom","We treated it as a funnel when it is also a cost."),
              ("A number we used","The 20 to 30 per cent taproom conversion has never been measured."),
              ("The census","Fifty self selected people cannot carry a strategy on their own."),
              ("Your marketing lead","We called Harry by the wrong name in a hundred and sixty places.")],
           top=3.1,col=3.4)

# 06
s_divider(prs,"Part 02","What to do","Three moves, ranked. And one we have ranked below them.")

# 07
s_content(prs,"Move 01",
  "Buy the databases of people who\nalready own a system.",
  "Harvey Norman lists six kegerator models plus a tap unit. Keg Land, Kegmaster and BenchTop "
  "hold buyer lists of their own. Every name on those lists owns hardware, so the price of the "
  "system, the most expensive objection your marketing fights, is already cleared. At about "
  "seventy dollars gross a keg, the data pays for itself quickly.")

# 08
s_content(prs,"Move 02",
  "Do a deal with Harvey Norman\ndirectly.",
  "A QR code at the point of sale. A referral rebate on each keg. Or a free keg with purchase "
  "if the liquor rules allow it. Monk is one of their closest allies and holds point contacts. "
  "This is a conversation, not a campaign, and it is worth having before any budget moves.")

# 09
s_content(prs,"Move 03",
  "Go at the tour operators, and run\nwholesale properly.",
  "Hop On, Urban Legends and Pineapple Tours all let the customer choose which breweries to "
  "visit. So your write up decides whether you get picked, and yours does not mention Midnight "
  "in Tokyo while Balter gets the imagery. On wholesale, one commercial account is worth about "
  "ten households, and it currently has no structure, no targets and nobody carrying it.")

# 10
s_content(prs,"And one we have ranked below them",
  "Pulling people into the taproom\nthrough social media.",
  "That was close to the centre of the plan we wrote for you. On this read it should not be. "
  "If the taproom is the half that loses money, spending the marketing budget filling it is the "
  "wrong instinct, however good the content is.")

# 11
s=s_content(prs,"What we need from you","Six decisions. None of them need research.",None)
rule_table(s,[("The taproom","Tasting and event venue, or keep paying for a drop in."),
              ("The range","Six core beers, or twenty seven taps."),
              ("Membership price","Two posters, two prices, both live in your venue today."),
              ("Active customer","Forty five, seventy five or ninety days. It decides every number we report."),
              ("Which plan is live","Three of your documents name three different core metrics."),
              ("Model A or B","Sell more systems, or convert the 116 who bought elsewhere.")],
           top=2.75,col=3.4,step=0.68)

# 12
s=s_content(prs,"How this splits","Harry takes these. We would do these for you.",None)
rule_table(s,[("Harry","Tour operator outreach. Fresh copy, photos, reviews and the award to each "
                       "operator. The referral reward nobody knows about."),
              ("Justin","The Harvey Norman conversation. Wholesale structure and targets. The six "
                        "decisions above."),
              ("Jewell","Database acquisition and the campaign against it. The Urban Legends booking "
                        "tool. Tracking, so you can see what any of it did.")],
           top=3.2,col=2.6,size=13)

# 13
s_closer(prs,"Decisions back by Friday 4 September. We will scope the three moves against them.",
             "Clent Jewell · clent@jewellprojects.com · Jewell Projects")

out="JP_TapThat_WhereYouAre_v01.pptx"
prs.save(out); print(f"{out}: {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
